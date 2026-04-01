"""Stakeholder presentation for Prompt2Test platform."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Dark theme palette --
BG       = RGBColor(0x0B, 0x11, 0x20)
CARD     = RGBColor(0x13, 0x1C, 0x33)
CARD2    = RGBColor(0x1A, 0x25, 0x42)
BLUE     = RGBColor(0x38, 0xBD, 0xF8)
PURPLE   = RGBColor(0xA7, 0x8B, 0xFA)
GREEN    = RGBColor(0x4A, 0xDE, 0x80)
ORANGE   = RGBColor(0xFB, 0x92, 0x3C)
RED      = RGBColor(0xF8, 0x71, 0x71)
YELLOW   = RGBColor(0xFB, 0xD3, 0x4D)
TEAL     = RGBColor(0x2D, 0xD4, 0xBF)
PINK     = RGBColor(0xF4, 0x72, 0xB6)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LG       = RGBColor(0xCB, 0xD5, 0xE1)
MG       = RGBColor(0x94, 0xA3, 0xB8)
DG       = RGBColor(0x47, 0x55, 0x69)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──
def _bg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG

def _t(s, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = "Segoe UI"; p.alignment = a
    return tf

def _rect(s, l, t, w, h, f=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = f; sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def _box(s, l, t, w, h, f=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = f; sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def _bar(s, l, t, w, c=BLUE):
    return _box(s, l, t, w, Pt(4), c)

def _pill(s, l, t, w, h, text, fc=BLUE, tc=WHITE, sz=14):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc; sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = tc
    p.font.bold = True; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER

def _circ(s, l, t, d, fc=BLUE, text="", tc=WHITE, sz=16):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc; sh.line.fill.background()
    if text:
        tf = sh.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = tc
        p.font.bold = True; p.alignment = PP_ALIGN.CENTER

def _arrow_r(s, l, t, w, c=BLUE):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, Inches(0.35))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background(); sh.shadow.inherit = False

def _arrow_d(s, l, t, h, c=BLUE):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, Inches(0.35), h)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background(); sh.shadow.inherit = False

def _title(s, text, sub="", ac=BLUE):
    _bg(s)
    _box(s, Inches(0), Inches(0), prs.slide_width, Pt(4), ac)
    _t(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), text, sz=32, c=WHITE, b=True)
    _bar(s, Inches(0.8), Inches(1.0), Inches(1.8), ac)
    if sub:
        _t(s, Inches(0.8), Inches(1.15), Inches(11), Inches(0.4), sub, sz=15, c=MG)


# ================================================================
# SLIDE 1 -- Title
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
_circ(s, Inches(-1), Inches(-1), Inches(4), RGBColor(0x10, 0x18, 0x2E))
_circ(s, Inches(10.5), Inches(5), Inches(4), RGBColor(0x10, 0x18, 0x2E))
_bar(s, Inches(0), Inches(0), prs.slide_width, PURPLE)

_t(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.5),
   "PROMPT2TEST", sz=18, c=PURPLE, b=True, a=PP_ALIGN.CENTER)
_t(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2),
   "AI-Powered Test Automation Platform", sz=48, c=WHITE, b=True, a=PP_ALIGN.CENTER)
_bar(s, Inches(5.5), Inches(3.2), Inches(2.3), PURPLE)
_t(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
   "Describe tests in plain English. AI generates, executes, and reports.\nBuilt on AWS in 15 days by 1 engineer.", sz=18, c=MG, a=PP_ALIGN.CENTER)

# Tech pills
techs = [("AWS Bedrock", BLUE), ("Claude Sonnet", PURPLE), ("Playwright", GREEN),
         ("React", BLUE), ("ECS Fargate", ORANGE), ("Aurora", TEAL)]
for i, (label, clr) in enumerate(techs):
    x = Inches(1.5) + Inches(i * 1.8)
    _pill(s, x, Inches(5.0), Inches(1.6), Inches(0.45), label, clr, WHITE, 12)

_t(s, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.4),
   "Stakeholder Review  |  March 2026", sz=13, c=DG, a=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 2 -- The Problem
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "The Problem", "Manual testing is slow, expensive, and doesn't scale", RED)

problems = [
    ("Manual test authoring takes hours",
     "QA engineers spend 60-70% of time writing test scripts instead of finding bugs.",
     "2-4 hours per test case", RED),
    ("Test maintenance is a constant burden",
     "Every UI change breaks existing tests. Teams spend more time fixing tests than writing new ones.",
     "40% of QA effort on maintenance", ORANGE),
    ("No live visibility during execution",
     "Tests run as black boxes. When they fail, engineers debug without seeing what happened.",
     "30-60 min to diagnose a failure", YELLOW),
    ("Scaling requires more headcount",
     "More features = more tests = more QA engineers. Linear cost growth.",
     "1 QA per 3-4 developers", PINK),
]

for i, (title, desc, stat, clr) in enumerate(problems):
    y = Inches(1.7) + Inches(i * 1.35)
    _rect(s, Inches(0.4), y, Inches(8.5), Inches(1.15), CARD)
    _box(s, Inches(0.4), y, Pt(4), Inches(1.15), clr)
    _circ(s, Inches(0.7), y + Inches(0.3), Inches(0.45), clr, str(i+1), WHITE, 16)
    _t(s, Inches(1.3), y + Inches(0.1), Inches(7.2), Inches(0.3), title, sz=16, c=WHITE, b=True)
    _t(s, Inches(1.3), y + Inches(0.45), Inches(7.2), Inches(0.5), desc, sz=12, c=LG)

    # Stat badge on right
    _rect(s, Inches(9.2), y + Inches(0.15), Inches(3.8), Inches(0.8), CARD2)
    _t(s, Inches(9.2), y + Inches(0.2), Inches(3.8), Inches(0.7), stat, sz=14, c=clr, b=True, a=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 3 -- The Solution
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "The Solution: Prompt2Test", "Describe what to test in plain English. AI does the rest.", GREEN)

# End-to-end flow
flow = [
    ("QA Engineer", "types a prompt", BLUE),
    ("React UI", "Amplify hosted", PURPLE),
    ("Bedrock AgentCore", "Claude Sonnet 4.5", GREEN),
    ("AI generates plan", "structured steps", ORANGE),
    ("ECS Fargate", "spins up browser", TEAL),
    ("Playwright MCP", "executes test", PINK),
    ("Live view", "noVNC stream", BLUE),
    ("PASS / FAIL", "saved to Aurora", GREEN),
]

for i, (title, sub, clr) in enumerate(flow):
    x = Inches(0.3) + Inches(i * 1.6)
    _rect(s, x, Inches(2.0), Inches(1.45), Inches(1.2), CARD)
    _bar(s, x, Inches(2.0), Inches(1.45), clr)
    _t(s, x + Inches(0.1), Inches(2.15), Inches(1.25), Inches(0.35), title, sz=11, c=clr, b=True, a=PP_ALIGN.CENTER)
    _t(s, x + Inches(0.1), Inches(2.55), Inches(1.25), Inches(0.3), sub, sz=9, c=MG, a=PP_ALIGN.CENTER)
    if i < 7:
        _t(s, x + Inches(1.4), Inches(2.35), Inches(0.3), Inches(0.3), ">", sz=14, c=DG, a=PP_ALIGN.CENTER)

# Three value props
props = [
    ("10x Faster Test Creation", "Describe in English, get structured test plan in seconds. No scripting.",
     "Minutes, not hours", GREEN),
    ("Live Browser Visibility", "Watch tests execute in real-time via noVNC. See exactly what the AI does.",
     "Zero black boxes", BLUE),
    ("Zero Infrastructure Overhead", "On-demand ECS tasks. Pay only when tests run. No always-on servers.",
     "$0.50/mo idle cost", ORANGE),
]

for i, (title, desc, stat, clr) in enumerate(props):
    x = Inches(0.3) + Inches(i * 4.3)
    _rect(s, x, Inches(3.8), Inches(4.1), Inches(2.8), CARD)
    _bar(s, x, Inches(3.8), Inches(4.1), clr)
    _t(s, x + Inches(0.2), Inches(3.95), Inches(3.7), Inches(0.35), title, sz=16, c=clr, b=True)
    _t(s, x + Inches(0.2), Inches(4.4), Inches(3.7), Inches(0.8), desc, sz=12, c=LG)
    _pill(s, x + Inches(0.2), Inches(5.5), Inches(2.5), Inches(0.35), stat, CARD2, clr, 11)

# Bottom callout
_rect(s, Inches(2.5), Inches(6.8), Inches(8.3), Inches(0.5), CARD2)
_t(s, Inches(2.5), Inches(6.82), Inches(8.3), Inches(0.45),
   "Result: QA engineers focus on WHAT to test, not HOW to test it.",
   sz=14, c=YELLOW, a=PP_ALIGN.CENTER, b=True)


# ================================================================
# SLIDE 4 -- What Was Built (4 screens)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "What Was Built", "4 fully functional screens -- all deployed and live", PURPLE)

screens = [
    ("Agent", PURPLE,
     "AI-powered test authoring",
     ["Pick a service from config", "Describe test in plain English",
      "AI generates structured plan", "Review & refine steps",
      "Execute with live browser view", "Save as reusable test case"]),
    ("Inventory", BLUE,
     "Test case management",
     ["Browse by env & service", "Semantic search (pgvector)",
      "View last run result", "Re-run / Replay tests",
      "Edit or delete test cases", "Run records history"]),
    ("Config", TEAL,
     "Service & account config",
     ["Services + URLs per env", "Test accounts per env/team",
      "DynamoDB backed", "Used by AI agent at runtime"]),
    ("Admin", ORANGE,
     "User & team management",
     ["Create / delete Cognito users", "Manage team groups",
      "Assign users to teams", "Admin-only visibility"]),
]

for i, (name, clr, desc, features) in enumerate(screens):
    x = Inches(0.2) + Inches(i * 3.3)
    _rect(s, x, Inches(1.6), Inches(3.1), Inches(5.6), CARD)
    _bar(s, x, Inches(1.6), Inches(3.1), clr)
    _t(s, x + Inches(0.2), Inches(1.75), Inches(2.7), Inches(0.35), name, sz=20, c=clr, b=True)
    _t(s, x + Inches(0.2), Inches(2.15), Inches(2.7), Inches(0.3), desc, sz=11, c=MG)

    for j, feat in enumerate(features):
        y = Inches(2.7) + Inches(j * 0.45)
        _circ(s, x + Inches(0.2), y + Inches(0.04), Inches(0.22), clr, str(j+1), WHITE, 8)
        _t(s, x + Inches(0.5), y, Inches(2.4), Inches(0.3), feat, sz=11, c=LG)

    _pill(s, x + Inches(0.2), Inches(5.7), Inches(2.7), Inches(0.35), "Deployed & Live", GREEN, WHITE, 11)


# ================================================================
# SLIDE 5 -- Feature Coverage
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "Feature Coverage", "15 features complete, 3 planned for next phase", GREEN)

features = [
    ("AI Test Plan Generation", "Complete", GREEN),
    ("Live Browser Automation", "Complete", GREEN),
    ("Live Browser View (noVNC)", "Complete", GREEN),
    ("Test Replay", "Complete", GREEN),
    ("Test Case Inventory", "Complete", GREEN),
    ("Semantic Search (pgvector)", "Complete", GREEN),
    ("Service Config Management", "Complete", GREEN),
    ("Test Account Management", "Complete", GREEN),
    ("Multi-Team Support", "Complete", GREEN),
    ("Multi-Environment (4 envs)", "Complete", GREEN),
    ("Admin Panel", "Complete", GREEN),
    ("Auth & Access Control", "Complete", GREEN),
    ("CI/CD -- 4 Pipelines", "Complete", GREEN),
    ("Exploratory Mode", "Complete", GREEN),
    ("Run Records", "Complete", GREEN),
    ("Test Metrics Dashboard", "Planned", DG),
    ("Slack / Email Notifications", "Planned", DG),
    ("Scheduled Test Runs", "Planned", DG),
]

for i, (name, status, clr) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.3) + Inches(col * 4.3)
    y = Inches(1.6) + Inches(row * 0.52)
    _rect(s, x, y, Inches(4.1), Inches(0.42), CARD)
    _t(s, x + Inches(0.15), y + Inches(0.05), Inches(2.8), Inches(0.3), name, sz=11, c=WHITE if clr == GREEN else MG, b=True)
    _pill(s, x + Inches(3.0), y + Inches(0.06), Inches(0.95), Inches(0.28), status, clr, WHITE, 9)

# Summary bar
_rect(s, Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.5), CARD2)
_t(s, Inches(0.3), Inches(6.82), Inches(12.7), Inches(0.45),
   "15 / 18 features shipped  |  83% complete  |  3 planned for Phase 2",
   sz=14, c=GREEN, a=PP_ALIGN.CENTER, b=True)


# ================================================================
# SLIDE 6 -- Architecture (high level)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "Architecture Overview", "18 AWS services, fully serverless, zero always-on infrastructure", BLUE)

# Layered architecture
layers = [
    ("USER LAYER", "QA Engineer > React SPA (Amplify + CloudFront)", BLUE, Inches(1.6)),
    ("AUTH LAYER", "Cognito User Pool > Identity Pool > JWT > IAM Roles", PURPLE, Inches(2.4)),
    ("AI LAYER", "Bedrock AgentCore > Claude Sonnet 4.5 > Strands SDK > MCP Tools", GREEN, Inches(3.2)),
    ("COMPUTE LAYER", "ECS Fargate (on-demand) > Playwright + noVNC > 3 Lambda Functions", ORANGE, Inches(4.0)),
    ("DATA LAYER", "Aurora Serverless v2 + pgvector  |  DynamoDB Config  |  RDS Data API", TEAL, Inches(4.8)),
    ("CI/CD LAYER", "3 CodePipelines > CodeBuild > ECR > Amplify Auto-deploy", PINK, Inches(5.6)),
]

for name, desc, clr, y in layers:
    _rect(s, Inches(0.4), y, Inches(12.5), Inches(0.65), CARD)
    _box(s, Inches(0.4), y, Pt(4), Inches(0.65), clr)
    _pill(s, Inches(0.6), y + Inches(0.12), Inches(2.0), Inches(0.38), name, clr, WHITE, 10)
    _t(s, Inches(2.8), y + Inches(0.15), Inches(9.8), Inches(0.35), desc, sz=12, c=LG)

# Arrows between layers
for i in range(5):
    y = Inches(2.25) + Inches(i * 0.8)
    _arrow_d(s, Inches(6.5), y, Inches(0.2), DG)

# Key metrics at bottom
metrics = [
    ("18", "AWS Services"), ("4", "Environments"), ("5", "GitHub Repos"),
    ("3", "CodePipelines"), ("20+", "Concurrent Users"), ("$0.50", "Idle Cost/mo"),
]
for i, (val, lbl) in enumerate(metrics):
    x = Inches(0.4) + Inches(i * 2.15)
    _rect(s, x, Inches(6.5), Inches(2.0), Inches(0.8), CARD2)
    _t(s, x, Inches(6.5), Inches(2.0), Inches(0.45), val, sz=22, c=BLUE, b=True, a=PP_ALIGN.CENTER)
    _t(s, x, Inches(6.95), Inches(2.0), Inches(0.3), lbl, sz=10, c=MG, a=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 7 -- Cost & Scaling
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "Cost & Scaling Model", "Zero always-on compute. Pay only when tests run.", TEAL)

# Cost comparison
_rect(s, Inches(0.4), Inches(1.7), Inches(5.8), Inches(2.5), CARD)
_bar(s, Inches(0.4), Inches(1.7), Inches(5.8), TEAL)
_t(s, Inches(0.6), Inches(1.85), Inches(5.4), Inches(0.35), "MONTHLY COST", sz=16, c=TEAL, b=True)

costs = [
    ("Idle (no tests running)", "$0.50/mo", "Aurora auto-pause, zero ECS tasks", GREEN),
    ("Light usage (10 tests/day)", "$15-25/mo", "On-demand Fargate + Bedrock calls", BLUE),
    ("Heavy usage (50+ tests/day)", "$30-70/mo", "Scales linearly with usage", ORANGE),
]
for i, (scenario, cost, detail, clr) in enumerate(costs):
    y = Inches(2.4) + Inches(i * 0.55)
    _t(s, Inches(0.8), y, Inches(3.0), Inches(0.3), scenario, sz=12, c=WHITE, b=True)
    _pill(s, Inches(3.8), y, Inches(1.3), Inches(0.32), cost, clr, WHITE, 11)
    _t(s, Inches(5.2), y + Inches(0.02), Inches(1.0), Inches(0.3), detail, sz=9, c=MG)

# vs traditional
_rect(s, Inches(6.5), Inches(1.7), Inches(6.4), Inches(2.5), CARD)
_bar(s, Inches(6.5), Inches(1.7), Inches(6.4), RED)
_t(s, Inches(6.7), Inches(1.85), Inches(5.4), Inches(0.35), "VS TRADITIONAL APPROACH", sz=16, c=RED, b=True)

_t(s, Inches(6.9), Inches(2.4), Inches(5.8), Inches(0.3),
   "Always-on Selenium Grid: $500-2000/mo", sz=13, c=RED)
_t(s, Inches(6.9), Inches(2.8), Inches(5.8), Inches(0.3),
   "QA engineer salaries: $8000-15000/mo per person", sz=13, c=RED)
_t(s, Inches(6.9), Inches(3.2), Inches(5.8), Inches(0.3),
   "Test maintenance overhead: 40% of QA time", sz=13, c=ORANGE)

# Scaling model
_rect(s, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.8), CARD)
_t(s, Inches(0.6), Inches(4.6), Inches(11.8), Inches(0.35), "SCALING MODEL", sz=16, c=WHITE, b=True)

scale_items = [
    ("Zero always-on tasks", "ECS tasks spin up only when a test is running. Shut down immediately after.", TEAL),
    ("1 task per user session", "Each test gets an isolated Fargate task with its own browser. No cross-contamination.", GREEN),
    ("20+ concurrent capacity", "Fargate scales to 20+ simultaneous test sessions with no pre-provisioning.", BLUE),
    ("Aurora auto-pause", "Database pauses after 5 min idle. Resumes in ~30s on first query. Near-zero idle cost.", PURPLE),
    ("Bedrock pay-per-token", "Claude Sonnet charges only for tokens used. No reserved capacity needed.", ORANGE),
]
for i, (title, desc, clr) in enumerate(scale_items):
    x = Inches(0.6) + Inches((i % 3) * 4.15)
    y = Inches(5.1) + Inches((i // 3) * 1.0)
    _rect(s, x, y, Inches(3.95), Inches(0.85), CARD2)
    _box(s, x, y, Pt(3), Inches(0.85), clr)
    _t(s, x + Inches(0.15), y + Inches(0.05), Inches(3.6), Inches(0.25), title, sz=12, c=clr, b=True)
    _t(s, x + Inches(0.15), y + Inches(0.35), Inches(3.6), Inches(0.4), desc, sz=10, c=LG)


# ================================================================
# SLIDE 8 -- Build Timeline (the 15-day story)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "Built in 15 Days", "1 engineer, aggressive AI-assisted development (Claude Code)", YELLOW)

# Phase timeline
phases = [
    ("Phase 1", "Infrastructure", "VPC, Cognito, Aurora, DynamoDB,\nECS Fargate, CDK IaC", "3 days", BLUE),
    ("Phase 2", "Backend", "3 Lambda functions,\nPlaywright MCP container", "2 days", GREEN),
    ("Phase 3", "AI Agent", "AgentCore, plan/automate/replay\nmodes, MCP integration", "4 days", ORANGE),
    ("Phase 4", "Frontend", "4 React screens, Cognito auth,\nnoVNC, streaming events", "5 days", PURPLE),
    ("Phase 5", "CI/CD + Docs", "3 CodePipelines, Amplify,\narchitecture docs", "1 day", TEAL),
]

for i, (phase, title, desc, duration, clr) in enumerate(phases):
    x = Inches(0.2) + Inches(i * 2.6)
    _rect(s, x, Inches(1.7), Inches(2.45), Inches(3.0), CARD)
    _bar(s, x, Inches(1.7), Inches(2.45), clr)
    _pill(s, x + Inches(0.1), Inches(1.85), Inches(1.0), Inches(0.3), phase, clr, WHITE, 10)
    _pill(s, x + Inches(1.2), Inches(1.85), Inches(1.1), Inches(0.3), duration, CARD2, clr, 10)
    _t(s, x + Inches(0.1), Inches(2.3), Inches(2.25), Inches(0.3), title, sz=15, c=WHITE, b=True)
    _t(s, x + Inches(0.1), Inches(2.7), Inches(2.25), Inches(0.8), desc, sz=11, c=LG)
    if i < 4:
        _t(s, x + Inches(2.35), Inches(2.9), Inches(0.3), Inches(0.3), ">", sz=16, c=DG)

# AI acceleration comparison
_rect(s, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.3), CARD)
_t(s, Inches(0.6), Inches(5.1), Inches(5.0), Inches(0.3), "AI ACCELERATION COMPARISON", sz=14, c=WHITE, b=True)

comparisons = [
    ("Traditional (no AI)", "~6 months", "4-5 person team", RED),
    ("Moderate AI", "~2 months", "2 person team", ORANGE),
    ("Aggressive AI (Claude Code)", "15 days", "1 engineer", GREEN),
]
for i, (approach, time, team, clr) in enumerate(comparisons):
    x = Inches(0.6) + Inches(i * 4.15)
    _rect(s, x, Inches(5.5), Inches(3.95), Inches(1.5), CARD2)
    _bar(s, x, Inches(5.5), Inches(3.95), clr)
    _t(s, x + Inches(0.2), Inches(5.65), Inches(3.5), Inches(0.3), approach, sz=13, c=clr, b=True)
    _t(s, x + Inches(0.2), Inches(6.05), Inches(3.5), Inches(0.4), time, sz=24, c=clr, b=True)
    _t(s, x + Inches(0.2), Inches(6.5), Inches(3.5), Inches(0.3), team, sz=11, c=MG)

# 12x callout
_pill(s, Inches(10.5), Inches(5.1), Inches(2.2), Inches(0.35), "12x faster", GREEN, WHITE, 13)


# ================================================================
# SLIDE 9 -- Key Architecture Decisions
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "Key Architecture Decisions", "Why we chose what we chose", ORANGE)

decisions = [
    ("DynamoDB for Config", "Not SSM Parameter Store",
     "Multi-team scoping with pk/sk, browser-direct access via Cognito Identity Pool, no Lambda proxy needed.", BLUE),
    ("Cognito Groups for Teams", "Not custom RBAC",
     "Zero-code team isolation. JWT carries group, all services scope data by team automatically.", PURPLE),
    ("ECS Fargate On-demand", "Not warm pool or EC2",
     "Zero idle cost. 60s cold start is acceptable. Was paying $216/mo for 3-task warm pool.", GREEN),
    ("Bedrock AgentCore", "Not self-hosted agent",
     "Managed container runtime. No ECS/EKS for the agent. Built-in scaling, versioning, and monitoring.", ORANGE),
    ("Aurora + pgvector", "Not OpenSearch or Pinecone",
     "Single database for relational data AND vector search. No extra service. RDS Data API for serverless access.", TEAL),
    ("CDK Infrastructure as Code", "Not console or CloudFormation",
     "TypeScript CDK for all 18 services. One-command deploy. Reproducible across accounts.", PINK),
]

for i, (title, alt, reason, clr) in enumerate(decisions):
    col = i % 2
    row = i // 2
    x = Inches(0.3) + Inches(col * 6.5)
    y = Inches(1.6) + Inches(row * 1.85)
    _rect(s, x, y, Inches(6.3), Inches(1.65), CARD)
    _box(s, x, y, Pt(4), Inches(1.65), clr)
    _t(s, x + Inches(0.2), y + Inches(0.1), Inches(3.5), Inches(0.3), title, sz=15, c=clr, b=True)
    _pill(s, x + Inches(3.8), y + Inches(0.1), Inches(2.3), Inches(0.28), alt, CARD2, MG, 9)
    _t(s, x + Inches(0.2), y + Inches(0.5), Inches(5.9), Inches(0.9), reason, sz=11, c=LG)


# ================================================================
# SLIDE 10 -- What's Next
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_title(s, "What's Next", "Phase 2 roadmap and expansion opportunities", YELLOW)

# Phase 2 items
_rect(s, Inches(0.4), Inches(1.6), Inches(6.0), Inches(3.5), CARD)
_bar(s, Inches(0.4), Inches(1.6), Inches(6.0), GREEN)
_t(s, Inches(0.6), Inches(1.75), Inches(5.5), Inches(0.35), "PHASE 2 -- PLANNED FEATURES", sz=16, c=GREEN, b=True)

phase2 = [
    ("Test Metrics Dashboard", "Aggregate pass/fail rates, trends, most-run services", GREEN),
    ("Slack / Email Notifications", "Alert on test failure or completion via webhook or SES", BLUE),
    ("Scheduled Test Runs", "Cron-based smoke tests -- nightly, per-deploy, or on-demand", ORANGE),
    ("Multi-Agent Collaboration", "Specialized agents for different test types (API, visual, perf)", PURPLE),
    ("Test Coverage Insights", "Map test cases to features/stories, identify coverage gaps", TEAL),
]
for i, (title, desc, clr) in enumerate(phase2):
    y = Inches(2.3) + Inches(i * 0.52)
    _circ(s, Inches(0.7), y + Inches(0.04), Inches(0.3), clr, str(i+1), WHITE, 10)
    _t(s, Inches(1.1), y, Inches(2.5), Inches(0.25), title, sz=12, c=clr, b=True)
    _t(s, Inches(3.6), y, Inches(2.6), Inches(0.4), desc, sz=10, c=LG)

# Expansion opportunities
_rect(s, Inches(6.7), Inches(1.6), Inches(6.2), Inches(3.5), CARD)
_bar(s, Inches(6.7), Inches(1.6), Inches(6.2), PURPLE)
_t(s, Inches(6.9), Inches(1.75), Inches(5.7), Inches(0.35), "EXPANSION OPPORTUNITIES", sz=16, c=PURPLE, b=True)

expansions = [
    ("API Testing", "Extend beyond browser tests to REST/GraphQL API automation", BLUE),
    ("Visual Regression", "Screenshot comparison to catch unintended UI changes", ORANGE),
    ("Performance Testing", "Load test generation from natural language descriptions", TEAL),
    ("Cross-team Rollout", "Onboard additional teams -- config is already multi-tenant", GREEN),
    ("CI/CD Integration", "Trigger Prompt2Test on every deployment automatically", PINK),
]
for i, (title, desc, clr) in enumerate(expansions):
    y = Inches(2.3) + Inches(i * 0.52)
    _circ(s, Inches(7.0), y + Inches(0.04), Inches(0.3), clr, str(i+1), WHITE, 10)
    _t(s, Inches(7.4), y, Inches(2.5), Inches(0.25), title, sz=12, c=clr, b=True)
    _t(s, Inches(9.9), y, Inches(2.8), Inches(0.4), desc, sz=10, c=LG)

# Key ask
_rect(s, Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.8), CARD)
_bar(s, Inches(0.4), Inches(5.4), Inches(12.5), YELLOW)
_t(s, Inches(0.6), Inches(5.55), Inches(11.8), Inches(0.35), "KEY TAKEAWAY", sz=16, c=YELLOW, b=True)
_t(s, Inches(0.6), Inches(6.0), Inches(11.8), Inches(0.5),
   "Prompt2Test is production-ready today. 15 features shipped, 18 AWS services integrated, multi-team/multi-env.",
   sz=14, c=LG)
_t(s, Inches(0.6), Inches(6.5), Inches(11.8), Inches(0.5),
   "Built by 1 engineer in 15 days using AI-assisted development -- demonstrating 12x acceleration over traditional approaches.",
   sz=14, c=GREEN, b=True)


# -- Save --
out = r"c:\MyProjects\AWS\Prompt2TestInfrastructure\documentation\Prompt2Test_Stakeholder_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
